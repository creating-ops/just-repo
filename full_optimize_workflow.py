#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT完整优化流程：一键执行所有优化步骤
整合以下脚本功能：
1. convert_ppt_to_pptx.py - PPT转PPTX
2. optimize_ppt_final.py - 基础优化
3. optimize_central_bank_red.py - 央行红配色
4. optimize_content_background.py - 内容背景优化（修复logo保留）
5. fix_toc_page.py - 目录页修复

输出：2026货币金银重点工作思路4.2_完整美化版.pptx
"""

import os
import sys
import shutil
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

# ===================== 全局配置 =====================

BASE_DIR = Path(r'D:\文档\report\report\just-repo')
INPUT_PPT = BASE_DIR / '2026货币金银重点工作思路4.2.ppt'
OUTPUT_FINAL = BASE_DIR / '2026货币金银重点工作思路4.2_完整美化版.pptx'

# 央行红配色
CENTRAL_BANK_RED = RGBColor(0xC4, 0x1E, 0x3A)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF8, 0xF8)

# ===================== 步骤函数 =====================

def step1_convert_ppt_to_pptx():
    """步骤1: PPT转PPTX"""
    print("\n" + "=" * 60)
    print("步骤1: PPT转PPTX")
    print("=" * 60)

    # 检查是否已有PPTX版本
    temp_pptx = BASE_DIR / '2026货币金银重点工作思路4.2.pptx'
    if temp_pptx.exists():
        print(f"已存在PPTX文件: {temp_pptx}")
        return str(temp_pptx)

    # 使用COM转换
    try:
        import win32com.client
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1

        deck = powerpoint.Presentations.Open(str(INPUT_PPT))
        deck.SaveAs(str(temp_pptx), 24)  # 24 = pptx格式
        deck.Close()
        powerpoint.Quit()

        print(f"✅ 转换成功: {temp_pptx}")
        return str(temp_pptx)
    except Exception as e:
        print(f"❌ 转换失败: {e}")
        return None

def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def move_to_bottom(slide, shape):
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def step2_basic_optimize(prs):
    """步骤2: 基础优化（标题样式统一等）"""
    print("\n" + "=" * 60)
    print("步骤2: 基础样式优化")
    print("=" * 60)

    # 遍历所有幻灯片，统一样式
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # 统一字体
                        if run.font.name and 'MiSans' in str(run.font.name):
                            run.font.name = 'Microsoft YaHei'

    print("✅ 基础优化完成")
    return prs

def step3_apply_central_bank_red(prs):
    """步骤3: 应用央行红配色（保留logo）"""
    print("\n" + "=" * 60)
    print("步骤3: 应用央行红配色")
    print("=" * 60)

    for idx, slide in enumerate(prs.slides, 1):
        print(f"  处理第 {idx}/{len(prs.slides)} 页...", end="")

        if idx == 1:  # 封面
            # 添加央行红背景（不移除任何元素）
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.33), Inches(7.5)
            )
            set_shape_fill(bg, CENTRAL_BANK_RED)
            bg.line.fill.background()
            move_to_bottom(slide, bg)

            # 添加金色装饰线
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(3), Inches(2.0),
                Inches(8.33), Pt(4)
            )
            set_shape_fill(line, GOLD)
            line.line.fill.background()

            # 更新文字颜色
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = ''.join([p.text for p in shape.text_frame.paragraphs])
                    if '2026货币金银' in text or '重点工作思路' in text:
                        for para in shape.text_frame.paragraphs:
                            para.alignment = PP_ALIGN.CENTER
                            for run in para.runs:
                                run.font.color.rgb = WHITE
                                run.font.size = Pt(44)
                                run.font.bold = True
                    elif '汇报人' in text:
                        for para in shape.text_frame.paragraphs:
                            para.alignment = PP_ALIGN.CENTER
                            for run in para.runs:
                                run.font.color.rgb = GOLD
                                run.font.size = Pt(18)

        elif idx == len(prs.slides):  # 结束页
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.33), Inches(7.5)
            )
            set_shape_fill(bg, CENTRAL_BANK_RED)
            bg.line.fill.background()
            move_to_bottom(slide, bg)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = ''.join([p.text for p in shape.text_frame.paragraphs])
                    if '谢谢' in text or '聆听' in text:
                        for para in shape.text_frame.paragraphs:
                            para.alignment = PP_ALIGN.CENTER
                            for run in para.runs:
                                run.font.color.rgb = WHITE
                                run.font.size = Pt(48)
        else:  # 内容页
            # 添加央行红标题栏
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.33), Inches(0.6)
            )
            set_shape_fill(header, CENTRAL_BANK_RED)
            header.line.fill.background()
            move_to_bottom(slide, header)

        print(" ✓")

    print("✅ 央行红配色应用完成（logo已保留）")
    return prs

def step4_optimize_content_background(prs):
    """步骤4: 优化内容页背景（保留logo）"""
    print("\n" + "=" * 60)
    print("步骤4: 优化内容页背景")
    print("=" * 60)

    for idx, slide in enumerate(prs.slides, 1):
        print(f"  处理第 {idx}/{len(prs.slides)} 页...", end="")

        if idx == 1 or idx == len(prs.slides):
            # 封面和结束页已在步骤3处理
            print("跳过")
            continue

        # 内容页：添加浅灰背景
        # 检查并移除右上角装饰图（但不移除logo）
        for shape in list(slide.shapes):
            try:
                if hasattr(shape, 'image') and hasattr(shape, 'top') and hasattr(shape, 'left'):
                    # 只移除右上角小装饰图（保留logo）
                    if shape.top < Inches(1.5) and shape.left > Inches(10) and shape.width < Inches(2):
                        sp = shape._element
                        sp.getparent().remove(sp)
            except:
                pass

        # 添加浅灰背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.33), Inches(7.5)
        )
        set_shape_fill(bg, LIGHT_GRAY)
        bg.line.fill.background()
        move_to_bottom(slide, bg)

        print(" ✓")

    print("✅ 内容页背景优化完成")
    return prs

def step5_fix_toc_page(prs):
    """步骤5: 修复目录页"""
    print("\n" + "=" * 60)
    print("步骤5: 修复目录页")
    print("=" * 60)

    toc_slide = prs.slides[1]

    # 修改棕褐色为央行红
    spTree = toc_slide.shapes._spTree
    for sp in spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sp'):
        for srgbClr in sp.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr'):
            val = srgbClr.get('val', '')
            if val in ['8B4513', 'A0522D', 'CD853F']:
                srgbClr.set('val', 'C41E3A')

    # 修复数字编号大小
    for shape in toc_slide.shapes:
        if shape.has_text_frame:
            text = ''.join([p.text for p in shape.text_frame.paragraphs]).strip()
            if text in ['01', '02', '03']:
                for para in shape.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(28)
                        run.font.color.rgb = WHITE
                        run.font.bold = True

    print("✅ 目录页修复完成")
    return prs

# ===================== 主执行 =====================

def main():
    """一键执行完整优化流程"""
    print("=" * 60)
    print("PPT完整美化流程")
    print("=" * 60)
    print(f"\n输入文件: {INPUT_PPT}")
    print(f"输出文件: {OUTPUT_FINAL}")

    # 步骤1: PPT转PPTX
    pptx_file = step1_convert_ppt_to_pptx()
    if not pptx_file:
        print("\n❌ 流程终止：PPT转换失败")
        return

    # 读取PPTX
    prs = Presentation(pptx_file)
    print(f"\n共 {len(prs.slides)} 页幻灯片")

    # 步骤2: 基础优化
    prs = step2_basic_optimize(prs)

    # 步骤3: 央行红配色
    prs = step3_apply_central_bank_red(prs)

    # 步骤4: 内容背景优化
    prs = step4_optimize_content_background(prs)

    # 步骤5: 目录页修复
    prs = step5_fix_toc_page(prs)

    # 保存最终文件
    print(f"\n保存最终文件: {OUTPUT_FINAL}")
    prs.save(str(OUTPUT_FINAL))

    # 验证logo是否保留
    final_prs = Presentation(str(OUTPUT_FINAL))
    logo_count = sum(1 for s in final_prs.slides[0].shapes if hasattr(s, 'image'))
    print(f"封面页图片数量: {logo_count} 张")

    print("\n" + "=" * 60)
    print("✅ 完整美化流程执行完毕！")
    print("=" * 60)
    print("\n执行顺序总结:")
    print("  1. convert_ppt_to_pptx.py - PPT转PPTX")
    print("  2. optimize_ppt_final.py - 基础样式优化")
    print("  3. optimize_central_bank_red.py - 央行红配色（保留logo）")
    print("  4. optimize_content_background.py - 内容背景优化（修复logo保留逻辑）")
    print("  5. fix_toc_page.py - 目录页修复")
    print(f"\n最终输出: {OUTPUT_FINAL}")
    print("=" * 60)

if __name__ == '__main__':
    main()